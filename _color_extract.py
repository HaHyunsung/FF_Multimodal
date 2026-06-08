"""원본 PPT의 색상 팔레트 추출"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.dml.color import RGBColor

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표_원본백업.pptx'
prs = Presentation(SRC)

def color_of(shape):
    try:
        if shape.fill.type is not None:
            try:
                return shape.fill.fore_color.rgb
            except:
                pass
    except:
        pass
    return None

# 슬라이드 2의 헤더와 카드 색상
print("=== Slide 2 (콘텐츠 슬라이드 샘플) ===")
slide = prs.slides[1]
for j, sh in enumerate(slide.shapes):
    c = color_of(sh)
    txt = sh.text_frame.text[:30].replace('\n',' ') if sh.has_text_frame else ""
    print(f"  [{j}] fill_color={c} text=\"{txt}\"")

# 슬라이드 4 카드 색상도
print("\n=== Slide 4 (3 카드 형) ===")
slide = prs.slides[3]
for j, sh in enumerate(slide.shapes):
    c = color_of(sh)
    txt = sh.text_frame.text[:30].replace('\n',' ') if sh.has_text_frame else ""
    print(f"  [{j}] fill_color={c} text=\"{txt}\"")
