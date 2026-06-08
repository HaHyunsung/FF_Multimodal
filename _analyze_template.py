"""
슬라이드 마스터와 각 슬라이드의 헤더/하단바 구조를 분석한다.
- 슬라이드 마스터에 있는 요소 (모든 슬라이드 공통)
- 각 슬라이드 고유의 요소 중 어떤 게 헤더성/번호바/하단바인지 식별
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표_원본백업.pptx'
prs = Presentation(SRC)

W, H = prs.slide_width, prs.slide_height
print(f"Slide size: {W} x {H} EMU = {W/914400:.1f}\" x {H/914400:.1f}\"")

# 슬라이드 마스터 분석
print(f"\n=== Slide Masters: {len(prs.slide_masters)} ===")
for mi, master in enumerate(prs.slide_masters):
    print(f"Master {mi}: {len(master.shapes)} shapes")
    for j, sh in enumerate(master.shapes):
        txt = sh.text_frame.text[:50] if sh.has_text_frame else "(no text)"
        print(f"  [{j}] {sh.shape_type} L={sh.left} T={sh.top} W={sh.width} H={sh.height} \"{txt}\"")

    print(f"  Slide Layouts: {len(master.slide_layouts)}")
    for li, layout in enumerate(master.slide_layouts):
        print(f"    Layout {li}: '{layout.name}' - {len(layout.shapes)} shapes")

# 슬라이드별 - 위치/크기 정보 (헤더 판별용)
print(f"\n=== Slides: {len(prs.slides)} ===")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (layout: {slide.slide_layout.name}) ---")
    for j, sh in enumerate(slide.shapes):
        L = sh.left if sh.left else 0
        T = sh.top if sh.top else 0
        W_ = sh.width if sh.width else 0
        H_ = sh.height if sh.height else 0
        txt = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text[:40].replace('\n', '|')
        # 위치를 인치로
        Li, Ti, Wi, Hi = L/914400, T/914400, W_/914400, H_/914400
        print(f"  [{j}] {sh.shape_type} pos=({Li:.2f},{Ti:.2f}) size=({Wi:.2f}x{Hi:.2f}) \"{txt}\"")
