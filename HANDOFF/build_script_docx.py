# -*- coding: utf-8 -*-
"""발표 스크립트(2장 이내) + 뒷장 QnA 참고자료 Word 생성."""
import os
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation

PPTX = "딥러닝응용_3조_최종발표.pptx"
OUT = "발표스크립트_및_QnA참고.docx"

prs = Presentation(PPTX)
titles = ["표지/인사", "왜 멀티모달인가", "프로젝트 목표", "데이터셋", "세 가지 모델",
          "시스템 파이프라인", "중간→최종 가설검증", "핵심 결과", "이미지 열화 견고성",
          "결함 유형별 성능", "결론", "한계와 향후", "마무리"]
notes = []
for sl in prs.slides:
    t = ""
    try:
        t = sl.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass
    notes.append(t)

doc = Document()
for s in doc.sections:
    s.top_margin = s.bottom_margin = Cm(1.2)
    s.left_margin = s.right_margin = Cm(1.4)


def set_font(style, name, size):
    style.font.name = name
    style.font.size = Pt(size)
    rpr = style.element.get_or_add_rPr()
    rf = rpr.find(qn('w:rFonts'))
    if rf is None:
        rf = OxmlElement('w:rFonts')
        rpr.append(rf)
    rf.set(qn('w:eastAsia'), name)
    rf.set(qn('w:ascii'), name)
    rf.set(qn('w:hAnsi'), name)


normal = doc.styles['Normal']
set_font(normal, '맑은 고딕', 9.5)
pf = normal.paragraph_format
pf.line_spacing_rule = WD_LINE_SPACING.SINGLE
pf.space_after = Pt(2)
pf.space_before = Pt(0)

NAVY = RGBColor(0x1E, 0x2A, 0x4A)


def head(txt, size=12):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(3)
    p.paragraph_format.space_after = Pt(3)
    r = p.add_run(txt)
    r.bold = True
    r.font.size = Pt(size)
    r.font.color.rgb = NAVY
    return p


def bullets(items):
    for it in items:
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Cm(0.4)
        p.add_run("· " + it)


def table(headers, rows):
    tb = doc.add_table(rows=1, cols=len(headers))
    tb.style = 'Table Grid'
    for j, h in enumerate(headers):
        c = tb.rows[0].cells[j]
        c.text = ""
        rr = c.paragraphs[0].add_run(h)
        rr.bold = True
        rr.font.size = Pt(9)
    for row in rows:
        cells = tb.add_row().cells
        for j, v in enumerate(row):
            cells[j].text = ""
            rr = cells[j].paragraphs[0].add_run(str(v))
            rr.font.size = Pt(9)
    return tb


# ===== Part 1: 발표 스크립트 =====
title = doc.add_paragraph()
tr = title.add_run("딥러닝응용 3조 — 최종 발표 스크립트")
tr.bold = True
tr.font.size = Pt(13)
tr.font.color.rgb = NAVY
title.paragraph_format.space_after = Pt(4)

for i, (t, n) in enumerate(zip(titles, notes)):
    if not n:
        n = ("(인사·팀 소개 — 안녕하세요, 딥러닝응용 3조 김주헌·하현성입니다.)"
             if i == 0 else "(생략)")
    p = doc.add_paragraph()
    rh = p.add_run("S%d. %s  " % (i + 1, t))
    rh.bold = True
    rh.font.color.rgb = NAVY
    p.add_run(n)

# ===== Part 2: QnA 참고자료 =====
doc.add_page_break()
title2 = doc.add_paragraph()
t2 = title2.add_run("[참고자료] 데이터·모델·결과 상세 (QnA 대비)")
t2.bold = True
t2.font.size = Pt(13)
t2.font.color.rgb = NAVY
title2.paragraph_format.space_after = Pt(4)

head("A. 데이터 핵심")
bullets([
    "출처: 사우스캐롤라이나대(USC) Future Factories — 실제 조립라인 모사 시설, 로봇 4대+컨베이어가 30시간 동안 로켓 모형 조립.",
    "규모: 총 166,001 측정 시점(1초 약 2회=1.95Hz), 285 사이클(제품 285개). 클래스 5종: Normal(약 54.7%), NoNose, NoNose_NoBody2, NoNose_NoBody2_NoBody1(3부품), NoBody1.",
    "두 모달리티: ① 센서 22채널은 '모든' 시점 기록. ② 카메라 2대 이미지는 부품이 보이는 특정 공정 단계(CycleState 4·9)에서 '사용'(촬영 자체는 전 구간).",
    "결함 생성: 의도적으로 부품을 빼고 조립 → 결함이 누적적(코→몸통2→몸통1 순). 그래서 3부품 결함 = 2부품 + 부품 하나 더.",
    "분류 단위: 'CycleState 단위'가 아니라 '매 측정 시점' 단위(매 줄마다 5클래스 중 하나). 같은 (사이클,state) 안에선 라벨이 같지만, 센서 입력(직전 50시점 윈도우)이 시점마다 달라 각각 별개 샘플로 평가.",
    "분할: cycle-wise 80/20(테스트 57사이클) — 같은 사이클이 train/test에 섞이지 않게 해 데이터 누수 방지. 지표=Weighted F1.",
    "평가셋 3종: (1) 센서 전체 30,132 시점, (2) cycle4·9 전체 9,521 시점=FAIR 비교셋(이미지 결측 59% 포함), (3) 그중 이미지 실제 존재 3,912 시점(41%)=단독 vs 융합 정면 비교 가능한 유일 구간.",
    "용어: '장'이 아니라 '시점'. 한 시점=카메라 2장. 9,521=cycle4·9 전 시점, 3,912=그중 이미지 있는 시점.",
])

head("B. 모델 핵심")
bullets([
    "센서 모델: BiLSTM(양방향 2층). 입력=직전 50시점×22채널 슬라이딩 윈도우 → 그 끝 시점의 상태 분류. (강의 10-2 RNN/LSTM)",
    "이미지 모델: ResNet18(ImageNet 전이학습), 카메라 2장 예측 평균. 학습률 1e-4(1e-3은 과대→학습 불안정). (강의 7-2 CNN/전이학습)",
    "멀티모달 융합 3종(아래 표). 이미지 결측 시 센서로 폴백. (강의 9-1 멀티모달)",
])
table(["슬라이드 표기", "한국어 설명", "하는 일"],
      [["Decision (-Fusion)", "고정 가중 융합", "두 모델의 '확률'을 고정 비율(w=0.55)로 평균"],
       ["Concat", "학습형(특징결합) 융합", "두 모델의 '특징'을 이어붙여 함께 학습"],
       ["Cross-Attention", "적응형 주의 융합", "센서가 이미지에 적응적으로 주의→비중 자동 조절"]])

head("C. 결과 수치 (iter6 기준)")
table(["평가셋", "Sensor", "Image", "Concat", "CrossAttn", "Decision"],
      [["이미지 존재 3,912", "0.907", "0.957", "0.977", "0.906", "0.957"],
       ["전체 9,521(결측포함)", "0.871", "0.826", "0.900", "0.871", "0.893"]])
bullets([
    "센서 단독 전체(30,132 시점): F1 0.9215.",
    "부트스트랩 유의성: (Decision−Sensor) 평균 +0.050, 95% 신뢰구간 [+0.039, +0.060] → 0보다 커서 통계적으로 유의. (1000회 재추출로 차이 분포 추정)",
    "결함 유형별(Concat, cycle4·9): NoNose 0.97 / 2부품(NoNose_NoBody2) 0.94 / 3부품 0.71(가장 어려움) / Normal 0.92.",
    "핵심: Concat이 두 평가셋 모두에서 최고. 'Decision이 결측 환경에서 더 낫다'는 근거 없음(전체셋도 Concat 0.900 ≥ Decision 0.893).",
])

head("D. 이미지 열화 견고성 실험 (학습된 모델 그대로, 테스트 이미지만 열화)")
table(["이미지 품질", "Image", "Concat", "Decision", "CrossAttn", "Sensor"],
      [["깨끗", "0.957", "0.977", "0.957", "0.906", "0.907"],
       ["중간열화", "0.681", "0.931", "0.775↓", "0.905", "0.907"],
       ["강열화", "0.028", "0.303↓", "0.116↓", "0.905", "0.907"]])
bullets([
    "↓=센서(0.907) 밑으로 추락. Decision(고정가중)이 가장 취약, Concat은 중간까지 버팀, Cross-Attention만 끝까지 견고.",
    "주의: 강열화(이미지 0.03)는 '거의 파괴' 수준의 극단 가정 — 경향 확인용.",
])

head("E. 예상 QnA")
qa = [
    ("매 시점 분류가 무슨 뜻?", "CycleState마다 1번이 아니라, state 4·9에 속한 '모든 측정 시점'을 각각 5클래스로 분류."),
    ("점수는 모든 사이클을 종합한 값?", "예. 57개 테스트 사이클의 해당 시점을 전부 모아 Weighted F1을 한 번 계산(사이클별 평균 아님). 헤드라인=cycle4·9 9,521(이미지 있는 건 3,912)."),
    ("3부품 결함 0.71이 왜 낮나?", "3부품=2부품+부품 하나라 시각·센서적으로 닮아 2부품과 혼동(재현율 0.66로 언더콜). 단순 양/불 이진이면 거의 완벽했을 것 — 난이도는 '유형 세분류'에서 옴. 향후 '부품별 멀티라벨/순서형' 접근이 대안."),
    ("이미지 약하면 노이즈로 끌어내린다던 중간발표와 모순 아닌가?", "중간발표 약세 원인은 데이터 불균형+단순 concat. iter6는 데이터 보강(2파트·2카메라)+융합 개선으로, 이미지가 약해도 융합이 센서 수준을 유지(끌려내려가지 않음). 즉 모델/데이터 개선의 결과."),
    ("부트스트랩 검정이 뭔가?", "테스트 표본을 1000번 복원추출해 (융합−센서) 차이의 분포를 만들고, 그 95% 구간이 0보다 크면 '우연이 아니다(유의)'로 판단."),
    ("왜 논문(NSF-MAP)과 다른 모델?", "논문은 EfficientNet+YOLO+지식주입으로 '다음 단계 예측'. 우리는 강의에서 배운 ResNet18·BiLSTM으로 '현재 상태 분류' — 과제 취지(배운 기법 적용)에 맞춤."),
    ("Concat 발음?", "'컨캣'(concatenation=특징 이어붙이기의 줄임말). Decision='디시전', Cross-Attention='크로스 어텐션'."),
    ("멀티모달이 정확도를 크게 안 올렸는데 의미가 있나?", "깨끗한 이미지가 이미 0.96으로 강해 정확도 이득은 작지만, 이미지가 나빠질 때 융합(특히 적응형)이 성능을 지킨다는 '견고성'을 정량 규명한 것이 기여."),
]
for q, a in qa:
    p = doc.add_paragraph()
    r = p.add_run("Q. " + q)
    r.bold = True
    p2 = doc.add_paragraph("A. " + a)
    p2.paragraph_format.left_indent = Cm(0.4)

doc.save(OUT)
print("저장 완료:", OUT)
