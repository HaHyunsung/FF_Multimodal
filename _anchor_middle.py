"""모든 슬라이드의 모든 텍스트 프레임을 수직 가운데 정렬"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'

prs = Presentation(SRC)

def process_shape(shape, count):
    """텍스트 프레임 있으면 vertical_anchor를 MIDDLE로"""
    if shape.has_text_frame:
        try:
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            count[0] += 1
        except Exception as e:
            pass
    # 그룹 안 도형 재귀
    if shape.shape_type == 6:  # GROUP
        try:
            for sub in shape.shapes:
                process_shape(sub, count)
        except:
            pass

count = [0]
for i, slide in enumerate(prs.slides, 1):
    before = count[0]
    for shape in slide.shapes:
        process_shape(shape, count)
    print(f"Slide {i}: {count[0] - before}개 텍스트 프레임 수정")

print(f"\n총 {count[0]}개 텍스트 프레임 → MIDDLE 정렬")
prs.save(SRC)
print("저장 완료.")
