"""
중간발표 PPT 빌드 v2 - 처음부터 새로 디자인
- 헤더(상단바, 번호원, 제목)만 유지
- 콘텐츠 영역은 발표 내용에 맞춰 완전히 새로 그림
- 색상: 1B2A4A(남색), 16A295(틸), E8584F(코랄), F29E16(오렌지), 93B64E(그린), 2676AC(블루)
"""
import sys, io, copy, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ────────────────────────────────────────────
# 색상 팔레트
# ────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x16, 0xA2, 0x95)
CORAL = RGBColor(0xE8, 0x58, 0x4F)
ORANGE = RGBColor(0xF2, 0x9E, 0x16)
GREEN = RGBColor(0x93, 0xB6, 0x4E)
BLUE = RGBColor(0x26, 0x76, 0xAC)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_BLUE = RGBColor(0xF0, 0xF4, 0xF8)
HIGHLIGHT = RGBColor(0xEB, 0xF5, 0xFB)
SOFT_PINK = RGBColor(0xFF, 0xF3, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_GRAY = RGBColor(0x9A, 0xA0, 0xA6)

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표_원본백업.pptx'
OUT = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'

# ────────────────────────────────────────────
# 헬퍼: 슬라이드에서 헤더 외 도형 모두 삭제
# ────────────────────────────────────────────
def remove_content_keep_header(slide, keep_indices):
    """keep_indices에 있는 shape만 남기고 나머지 모두 삭제"""
    # 역순으로 삭제 (인덱스 어긋남 방지)
    sp_tree = slide.shapes._spTree
    shapes_to_remove = []
    for i, sh in enumerate(slide.shapes):
        if i not in keep_indices:
            shapes_to_remove.append(sh)
    for sh in shapes_to_remove:
        sp_tree.remove(sh._element)


# ────────────────────────────────────────────
# 헬퍼: 도형 생성
# ────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=None, line_color=None,
             line_width=None, shadow=False):
    """둥근 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()

    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()

    # 모서리 라운드 작게 조정
    try:
        adj = shape.adjustments
        adj[0] = 0.05
    except:
        pass

    if not shadow:
        # 그림자 제거
        spPr = shape._element.spPr
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        from lxml import etree
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))

    return shape


def add_circle(slide, left, top, size, fill, text="", font_color=WHITE, font_size=14, bold=True):
    """원 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.color.rgb = font_color
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.name = "맑은 고딕"
    # 그림자 제거
    spPr = shape._element.spPr
    from lxml import etree
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return shape


def add_text(slide, left, top, width, height, text, font_size=14,
             font_color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.color.rgb = font_color
        r.font.bold = bold
        r.font.name = font_name
    return tb


def add_arrow(slide, left, top, width, height, color=NAVY):
    """오른쪽 화살표"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ────────────────────────────────────────────
# 슬라이드 제목 텍스트 변경 (헤더 보존)
# ────────────────────────────────────────────
def set_header_title(slide, new_title):
    """슬라이드의 헤더 제목(shape[2])만 텍스트 변경"""
    title_shape = list(slide.shapes)[2]
    if title_shape.has_text_frame:
        # 첫 paragraph 첫 run의 스타일 유지하면서 텍스트 교체
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = new_title
            for r in p.runs[1:]:
                r.text = ''
        else:
            p.text = new_title
        # 나머지 paragraph 비우기
        for para in tf.paragraphs[1:]:
            para.text = ''


# ────────────────────────────────────────────
# 발표 빌드 시작
# ────────────────────────────────────────────
print(f"Loading: {SRC}")
prs = Presentation(SRC)

# ====================================================================
# Slide 1: 표지
# ====================================================================
print("\n[Slide 1] 표지")
slide = prs.slides[0]
# 표지는 GROUP 안에 있는 텍스트를 변경
# GROUP 내부 텍스트 5개 (제목 3줄 + 부제 + 학교 + 조원)
group = list(slide.shapes)[1]  # GROUP

new_texts = [
    "멀티모달 딥러닝 기반\n제조 공정 이상 탐지",
    "이미지와 시계열 센서 데이터를 활용한\n단일 vs 멀티모달 모델 비교 실험",
    "아주대학교 · 딥러닝응용 · 기말 프로젝트",
    "[3조] 김주헌 · 하현성",
]
sub_shapes = [s for s in group.shapes if s.has_text_frame and s.text_frame.text.strip()]
for i, sh in enumerate(sub_shapes):
    if i < len(new_texts):
        tf = sh.text_frame
        lines = new_texts[i].split('\n')
        # 첫 paragraph의 첫 run 스타일 유지
        first_p = tf.paragraphs[0]
        if first_p.runs:
            first_p.runs[0].text = lines[0]
            for r in first_p.runs[1:]:
                r.text = ''
        else:
            first_p.text = lines[0]
        # 추가 라인
        existing = list(tf.paragraphs)
        for j, line in enumerate(lines[1:], start=1):
            if j < len(existing):
                if existing[j].runs:
                    existing[j].runs[0].text = line
                    for r in existing[j].runs[1:]:
                        r.text = ''
                else:
                    existing[j].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
        for k in range(len(lines), len(existing)):
            existing[k].text = ''


# ====================================================================
# Slide 2: 01 / 프로젝트 동기
# ====================================================================
print("[Slide 2] 01 프로젝트 동기")
slide = prs.slides[1]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "프로젝트 동기")

# 좌측 큰 인용구
add_rect(slide, 0.42, 1.4, 12.5, 0.9, fill=HIGHLIGHT)
add_text(slide, 0.6, 1.45, 12.2, 0.8,
         "현장에서 비전 단독 양/불 판정은 조명·각도·가려짐으로 한계가 명확하고,\n센서 단독 판정도 노이즈·고장으로 신뢰하기 어려운 상황이 자주 발생한다.",
         font_size=14, font_color=NAVY, bold=False)

# 3개 카드
card_y = 2.55
card_h = 3.0
cards = [
    (0.42, CORAL, "1", "현장 경험 기반 문제 인식",
     "팀원 전원 자동화 장비 업계 종사\n현업의 딥러닝 불량 판정 사례를 직접 접하며 비전 단독 방식의 한계를 체감"),
    (4.79, BLUE, "2", "두 모달리티의 상호 보완",
     "이미지: 시각적 결손은 잘 잡지만 가려짐·조명에 취약\n센서: 동작 패턴은 정확하나 시각 정보는 부재\n→ 융합 시 약점 보완 가능성"),
    (9.16, ORANGE, "3", "프로젝트의 명확한 목표",
     "이미지·센서·융합 3개 모델을 동일 조건에서 학습·평가\n→ 멀티모달 융합이 실제로 어떤 이점을 주는지 정량 검증"),
]
for x, color, num, title, body in cards:
    add_rect(slide, x, card_y, 3.75, card_h, fill=LIGHT_BG)
    add_circle(slide, x + 0.2, card_y + 0.18, 0.6, color, num, font_size=18)
    add_text(slide, x + 0.2, card_y + 0.95, 3.45, 0.55, title,
             font_size=14, font_color=NAVY, bold=True)
    add_text(slide, x + 0.2, card_y + 1.55, 3.45, 1.4, body,
             font_size=11, font_color=DARK_GRAY)

# 하단 결론 바
add_rect(slide, 0.42, 5.75, 12.5, 0.65, fill=NAVY)
add_text(slide, 0.6, 5.78, 12.2, 0.6,
         "→  멀티모달 융합이 단일 모달리티 대비 실제 어떤 성능 차이를 만드는지 직접 구현하고 비교한다",
         font_size=14, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 3: 02 / 프로젝트 목표
# ====================================================================
print("[Slide 3] 02 프로젝트 목표")
slide = prs.slides[2]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "프로젝트 목표 및 접근")

# 큰 부제
add_text(slide, 0.42, 1.3, 12.5, 0.5,
         "Future Factories Dataset 기반 3-모델 비교 프레임워크",
         font_size=18, font_color=NAVY, bold=True)

# 3개 모델 박스
model_y = 2.0
model_h = 3.4
models = [
    (0.42, BLUE, "Model 1", "센서 단독", "BiLSTM",
     ["입력: 22 channel × 50 step", "전체 사이클 활용 가능", "Sliding Window 시퀀스",
      "Class Weight 불균형 처리"]),
    (4.79, TEAL, "Model 2", "이미지 단독", "ResNet18",
     ["ImageNet 사전학습 활용", "Transfer Learning + Fine-tune", "Data Augmentation 적용",
      "Cycle 4·9 구간만 사용"]),
    (9.16, CORAL, "Model 3", "멀티모달 융합", "Decision-Level Fusion",
     ["BiLSTM 256 + ResNet18 512", "Concat → FC Classifier", "두 모달리티 상호 보완",
      "이미지 없는 시점은 센서만"]),
]
for x, color, tag, name, arch, bullets in models:
    add_rect(slide, x, model_y, 3.75, model_h, fill=LIGHT_BG)
    # 상단 컬러 바
    add_rect(slide, x, model_y, 3.75, 0.55, fill=color)
    add_text(slide, x + 0.2, model_y + 0.08, 3.45, 0.4, tag,
             font_size=12, font_color=WHITE, bold=True)
    # 모델 이름
    add_text(slide, x + 0.2, model_y + 0.75, 3.45, 0.45, name,
             font_size=18, font_color=NAVY, bold=True)
    # 아키텍처
    add_text(slide, x + 0.2, model_y + 1.25, 3.45, 0.4, arch,
             font_size=12, font_color=color, bold=True)
    # 불릿
    bullet_text = "\n".join(f"•  {b}" for b in bullets)
    add_text(slide, x + 0.2, model_y + 1.8, 3.45, 1.6, bullet_text,
             font_size=11, font_color=DARK_GRAY)

# 하단 결론
add_rect(slide, 0.42, 5.65, 12.5, 0.7, fill=HIGHLIGHT)
add_text(slide, 0.6, 5.7, 12.2, 0.65,
         "동일 데이터 · 동일 평가 지표 (Accuracy, Precision, Recall, F1, Confusion Matrix) 로 세 모델을 비교하여 멀티모달의 효과를 정량 분석",
         font_size=12, font_color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 4: 03 / 관련 연구 NSF-MAP
# ====================================================================
print("[Slide 4] 03 관련 연구")
slide = prs.slides[3]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "관련 연구: NSF-MAP (IJCAI 2025)")

# 좌측 큰 박스: 논문 요약
add_text(slide, 0.42, 1.35, 7.0, 0.4,
         "NSF-MAP의 단계적 발전 전략",
         font_size=16, font_color=NAVY, bold=True)

# 3단계 카드 (세로 정렬)
stages = [
    (BLUE, "P1", "Decision-Level Fusion", "Accuracy 72%",
     "EfficientNet-B0 + Autoencoder의 특징을 Concat"),
    (TEAL, "P2", "+ Transfer Learning", "Accuracy 88%",
     "Encoder 동결 → Decoder Fine-tune 으로 과적합 방지"),
    (ORANGE, "P3", "+ Knowledge Infusion", "Accuracy 93%",
     "Process Ontology의 센서 범위 위반에 패널티"),
]
y0 = 1.85
for i, (color, tag, title, acc, desc) in enumerate(stages):
    y = y0 + i * 1.15
    add_rect(slide, 0.42, y, 7.0, 1.0, fill=LIGHT_BG)
    # 좌측 컬러 바
    add_rect(slide, 0.42, y, 0.7, 1.0, fill=color)
    add_text(slide, 0.42, y + 0.32, 0.7, 0.4, tag,
             font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 본문
    add_text(slide, 1.25, y + 0.08, 4.5, 0.45, title,
             font_size=13, font_color=NAVY, bold=True)
    add_text(slide, 1.25, y + 0.50, 5.5, 0.45, desc,
             font_size=10, font_color=DARK_GRAY)
    # 우측 정확도
    add_text(slide, 5.7, y + 0.25, 1.65, 0.5, acc,
             font_size=14, font_color=color, bold=True, align=PP_ALIGN.RIGHT)

# 우측: 우리 프로젝트와의 차이
add_rect(slide, 7.7, 1.35, 5.2, 4.2, fill=HIGHLIGHT)
add_text(slide, 7.85, 1.45, 5.0, 0.4,
         "본 프로젝트와의 차이",
         font_size=14, font_color=NAVY, bold=True)

diff_text = (
    "·  시계열: Autoencoder → BiLSTM\n"
    "    (강의 학습 내용 직접 적용)\n\n"
    "·  이미지: EfficientNet-B0 → ResNet18\n"
    "    (단순화 및 학습 시간 단축)\n\n"
    "·  핵심 비교 대상\n"
    "    NSF-MAP: P1 → P2 → P3 (방법론 비교)\n"
    "    본 프로젝트: 센서 vs 이미지 vs 융합\n"
    "    (모달리티 효과 비교)\n\n"
    "·  Knowledge Infusion / Ontology는\n"
    "    여건상 향후 확장 과제로 분류"
)
add_text(slide, 7.85, 1.9, 5.0, 3.5, diff_text,
         font_size=10, font_color=DARK_GRAY)

# 하단 결론
add_rect(slide, 0.42, 5.7, 12.5, 0.6, fill=NAVY)
add_text(slide, 0.6, 5.72, 12.2, 0.55,
         "→  본 프로젝트는 NSF-MAP을 참고하되, 강의 학습 내용 범위 안에서 멀티모달 융합의 효과 자체를 검증하는 데 집중",
         font_size=12, font_color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 5: 04 / 데이터셋
# ====================================================================
print("[Slide 5] 04 데이터셋")
slide = prs.slides[4]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "데이터셋: Future Factories")

# 좌측: 데이터셋 정보
add_text(slide, 0.42, 1.35, 7.0, 0.45,
         "Future Factories Lab Dataset (USC, 2024)",
         font_size=15, font_color=NAVY, bold=True)

info_y = 1.95
info_items = [
    ("조립 라인", "산업용 로봇 4대 + 컨베이어 구성"),
    ("조립 대상", "4개 부품 로켓 모형 (의도적 결함 주입)"),
    ("수집 규모", "30시간 연속 · 166,001 레코드 · 285 사이클"),
    ("센서 데이터", "그리퍼 로드셀·포텐셔미터, 로봇 관절각도,\n컨베이어 속도·온도 등 40+ 채널"),
    ("이미지 데이터", "카메라 2대 동기화 촬영 (총 ~166K 장)"),
]
for i, (label, body) in enumerate(info_items):
    y = info_y + i * 0.6
    add_rect(slide, 0.42, y, 1.8, 0.5, fill=NAVY)
    add_text(slide, 0.42, y + 0.05, 1.8, 0.4, label,
             font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 2.35, y - 0.02, 5.0, 0.6, body,
             font_size=10, font_color=DARK_GRAY)

# 우측: 라벨 분포 표
add_text(slide, 7.7, 1.35, 5.2, 0.45,
         "라벨 분포 (소수 클래스 병합 후)",
         font_size=15, font_color=NAVY, bold=True)

# 표 그리기
labels_data = [
    ("Normal", "90,775", "54.7%", BLUE),
    ("NoNose, NoBody2, NoBody1", "26,628", "16.0%", ORANGE),
    ("NoNose, NoBody2", "25,206", "15.2%", ORANGE),
    ("NoNose", "19,307", "11.7%", ORANGE),
    ("NoBody1 (병합)", "4,016", "2.4%", CORAL),
]
table_y = 1.95
add_rect(slide, 7.7, table_y, 5.2, 0.4, fill=NAVY)
add_text(slide, 7.85, table_y + 0.05, 2.5, 0.3, "Class",
         font_size=10, font_color=WHITE, bold=True)
add_text(slide, 10.4, table_y + 0.05, 1.2, 0.3, "Count",
         font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
add_text(slide, 11.7, table_y + 0.05, 1.1, 0.3, "Ratio",
         font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

for i, (name, cnt, ratio, color) in enumerate(labels_data):
    y = table_y + 0.4 + i * 0.42
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, 7.7, y, 5.2, 0.4, fill=bg)
    # 컬러 점
    add_circle(slide, 7.78, y + 0.13, 0.13, color, "", font_size=8)
    add_text(slide, 8.0, y + 0.04, 2.3, 0.3, name,
             font_size=9, font_color=DARK_GRAY)
    add_text(slide, 10.3, y + 0.04, 1.3, 0.3, cnt,
             font_size=9, font_color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, 11.65, y + 0.04, 1.15, 0.3, ratio,
             font_size=9, font_color=DARK_GRAY, align=PP_ALIGN.RIGHT)

# 하단 박스: 처리 전략
add_rect(slide, 0.42, 5.55, 12.5, 0.85, fill=HIGHLIGHT)
add_text(slide, 0.6, 5.62, 4.5, 0.3, "전처리 전략",
         font_size=11, font_color=NAVY, bold=True)
add_text(slide, 0.6, 5.93, 12.2, 0.45,
         "Cycle-wise 80/20 split (data leakage 방지) · Sliding Window 시퀀스 (T=50) · StandardScaler 정규화 · Class Weight로 불균형 보정",
         font_size=10, font_color=DARK_GRAY)


# ====================================================================
# Slide 6: 05 / 우리의 아키텍처
# ====================================================================
print("[Slide 6] 05 아키텍처")
slide = prs.slides[5]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "전체 아키텍처")

# 다이어그램 그리기 - 박스 + 화살표
# 입력 박스 2개 (좌측)
add_rect(slide, 0.5, 1.85, 2.4, 1.2, fill=BLUE)
add_text(slide, 0.55, 1.95, 2.3, 0.4, "Sensor Data",
         font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.55, 2.35, 2.3, 0.6, "22 channels\n50 time steps",
         font_size=10, font_color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 4.3, 2.4, 1.2, fill=TEAL)
add_text(slide, 0.55, 4.4, 2.3, 0.4, "Image Data",
         font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.55, 4.8, 2.3, 0.6, "224 × 224 × 3\n(Cycle 4·9)",
         font_size=10, font_color=WHITE, align=PP_ALIGN.CENTER)

# 화살표 1
add_arrow(slide, 2.95, 2.30, 0.6, 0.3, NAVY)
add_arrow(slide, 2.95, 4.75, 0.6, 0.3, NAVY)

# 인코더 박스
add_rect(slide, 3.65, 1.85, 2.6, 1.2, fill=LIGHT_BG, line_color=BLUE, line_width=1.5)
add_text(slide, 3.7, 1.95, 2.5, 0.4, "BiLSTM",
         font_size=13, font_color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 3.7, 2.35, 2.5, 0.6, "hidden 128 × 2-layer\nbidirectional",
         font_size=10, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 3.65, 4.3, 2.6, 1.2, fill=LIGHT_BG, line_color=TEAL, line_width=1.5)
add_text(slide, 3.7, 4.4, 2.5, 0.4, "ResNet18",
         font_size=13, font_color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 3.7, 4.8, 2.5, 0.6, "ImageNet pretrained\nbackbone (frozen)",
         font_size=10, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 화살표 2
add_arrow(slide, 6.3, 2.30, 0.5, 0.3, NAVY)
add_arrow(slide, 6.3, 4.75, 0.5, 0.3, NAVY)

# 특징 벡터
add_rect(slide, 6.9, 1.85, 1.7, 1.2, fill=BLUE)
add_text(slide, 6.95, 1.95, 1.6, 0.4, "f_sensor",
         font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 6.95, 2.4, 1.6, 0.5, "(256)",
         font_size=14, font_color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 6.9, 4.3, 1.7, 1.2, fill=TEAL)
add_text(slide, 6.95, 4.4, 1.6, 0.4, "f_image",
         font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 6.95, 4.85, 1.6, 0.5, "(512)",
         font_size=14, font_color=WHITE, align=PP_ALIGN.CENTER)

# 화살표 3 (concat으로 합류)
add_arrow(slide, 8.65, 2.30, 0.4, 0.3, NAVY)
add_arrow(slide, 8.65, 4.75, 0.4, 0.3, NAVY)

# Concat 박스
add_rect(slide, 9.1, 2.9, 2.0, 1.5, fill=CORAL)
add_text(slide, 9.15, 3.0, 1.9, 0.4, "Concat",
         font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 9.15, 3.45, 1.9, 0.4, "[f_s ; f_i]",
         font_size=12, font_color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 9.15, 3.85, 1.9, 0.4, "(768)",
         font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 화살표 4
add_arrow(slide, 11.15, 3.55, 0.4, 0.3, NAVY)

# 출력 박스
add_rect(slide, 11.6, 2.9, 1.4, 1.5, fill=NAVY)
add_text(slide, 11.6, 3.0, 1.4, 0.4, "FC",
         font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 11.6, 3.45, 1.4, 0.4, "Dropout",
         font_size=10, font_color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 11.6, 3.9, 1.4, 0.5, "5-class",
         font_size=11, font_color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# 하단 설명
add_rect(slide, 0.42, 5.85, 12.5, 0.55, fill=HIGHLIGHT)
add_text(slide, 0.6, 5.9, 12.2, 0.5,
         "Decision-Level Fusion: 각 모달리티의 특징을 독립적으로 추출 후 Concatenate → 두 모달리티가 보완 가능",
         font_size=11, font_color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 7: 06 / 모델별 학습 전략
# ====================================================================
print("[Slide 7] 06 학습 전략")
slide = prs.slides[6]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "모델별 학습 전략")

models_detail = [
    (0.42, BLUE, "M1", "Sensor BiLSTM",
     [("입력", "22채널 × 50 step 시퀀스"),
      ("적용 범위", "전체 사이클 (~150K 시퀀스)"),
      ("핵심", "양방향 LSTM으로\n과거·미래 컨텍스트 모두 활용"),
      ("학습", "Class Weight\nReduceLROnPlateau, Early Stop")]),
    (4.79, TEAL, "M2", "Image ResNet18",
     [("입력", "224 × 224 RGB 이미지"),
      ("적용 범위", "Cycle 4·9 시점 (~9,500장)"),
      ("핵심", "ImageNet 사전학습 활용\nTransfer Learning"),
      ("학습", "Augmentation\n(Flip, Rotation, ColorJitter)")]),
    (9.16, CORAL, "M3", "Multimodal Fusion",
     [("입력", "센서 시퀀스 + 동기화 이미지"),
      ("적용 범위", "Cycle 4·9 (이미지 보유 시점)"),
      ("핵심", "Decision-Level Fusion\nResNet 동결 + LSTM 학습"),
      ("학습", "이미지 없는 샘플은\n마스킹 처리하여 견고성 확보")]),
]
for x, color, tag, title, rows in models_detail:
    add_rect(slide, x, 1.4, 3.75, 4.5, fill=LIGHT_BG)
    # 상단 컬러 바
    add_rect(slide, x, 1.4, 3.75, 0.7, fill=color)
    add_text(slide, x + 0.2, 1.48, 1.0, 0.4, tag,
             font_size=14, font_color=WHITE, bold=True)
    add_text(slide, x + 1.2, 1.5, 2.5, 0.4, title,
             font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    # 행
    for i, (label, body) in enumerate(rows):
        y = 2.3 + i * 0.85
        add_text(slide, x + 0.2, y, 3.45, 0.3, label,
                 font_size=10, font_color=color, bold=True)
        add_text(slide, x + 0.2, y + 0.28, 3.45, 0.55, body,
                 font_size=10, font_color=DARK_GRAY)

# 하단 박스
add_rect(slide, 0.42, 6.05, 12.5, 0.4, fill=NAVY)
add_text(slide, 0.6, 6.07, 12.2, 0.36,
         "동일 5-class · 동일 Cycle-wise split · 동일 평가 지표 → 모달리티 차이만 분리하여 비교",
         font_size=11, font_color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 8: 07 / 실행 환경 (Kaggle)
# ====================================================================
print("[Slide 8] 07 실행 환경")
slide = prs.slides[7]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "실행 환경 및 데이터 제약")

# 좌측: 제약 사항
add_rect(slide, 0.42, 1.35, 6.0, 0.55, fill=CORAL)
add_text(slide, 0.6, 1.43, 5.8, 0.4, "마주한 제약",
         font_size=14, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

issues = [
    ("이미지 원본 580GB", "Kaggle의 Multi-Modal Dataset 6 파트의 총 용량 (각 ~104GB)\n→ 로컬 PC 다운로드 사실상 불가능"),
    ("로컬 GPU 부재", "팀원 PC에 학습용 GPU 없음\n→ ResNet18 학습 시 CPU만으로는 비현실적"),
    ("데이터 동기화", "JSON + 이미지 폴더 구조\n시계열 + 이미지 시점 매칭 별도 필요"),
]
y0 = 2.0
for i, (title, body) in enumerate(issues):
    y = y0 + i * 1.25
    add_rect(slide, 0.42, y, 6.0, 1.15, fill=SOFT_PINK)
    add_text(slide, 0.55, y + 0.1, 5.8, 0.35, title,
             font_size=12, font_color=CORAL, bold=True)
    add_text(slide, 0.55, y + 0.48, 5.8, 0.7, body,
             font_size=10, font_color=DARK_GRAY)

# 우측: 해결책
add_rect(slide, 6.92, 1.35, 6.0, 0.55, fill=TEAL)
add_text(slide, 7.1, 1.43, 5.8, 0.4, "해결 전략",
         font_size=14, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

sols = [
    ("Kaggle Notebook 활용", "Kaggle 서버에 데이터 마운트 → 다운로드 없이 직접 사용\n무료 GPU (Tesla T4 × 2) 주 30시간 제공"),
    ("이미지 파트 1개만 사용", "BATCH 1000~30000 (전체 1/6) 만 로드\n53,881장 중 9,542장 매칭 (~18%, NSF-MAP은 ~15K 사용)"),
    ("재현성 및 결과 백업", "Random Seed 42 고정 (numpy·torch·random)\n학습 가중치(.pt) + 결과(JSON) Kaggle Output에 저장"),
]
for i, (title, body) in enumerate(sols):
    y = y0 + i * 1.25
    add_rect(slide, 6.92, y, 6.0, 1.15, fill=LIGHT_BG)
    add_text(slide, 7.05, y + 0.1, 5.8, 0.35, title,
             font_size=12, font_color=TEAL, bold=True)
    add_text(slide, 7.05, y + 0.48, 5.8, 0.7, body,
             font_size=10, font_color=DARK_GRAY)

# 하단 결론
add_rect(slide, 0.42, 5.95, 12.5, 0.5, fill=NAVY)
add_text(slide, 0.6, 5.98, 12.2, 0.45,
         "→  Kaggle 환경 채택으로 데이터·연산 제약을 동시에 해결하면서 학술적으로 유의미한 샘플 규모 확보",
         font_size=11, font_color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 9: 08 / 최종 결과
# ====================================================================
print("[Slide 9] 08 최종 결과")
slide = prs.slides[8]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "최종 실험 결과")

# 좌측: 최우수 모델 강조 (센서)
add_rect(slide, 0.42, 1.35, 4.5, 4.8, fill=NAVY)
add_text(slide, 0.6, 1.5, 4.3, 0.4, "Best Model — Sensor (BiLSTM)",
         font_size=12, font_color=TEAL, bold=True)
add_text(slide, 0.6, 1.95, 4.3, 0.45, "BiLSTM",
         font_size=20, font_color=WHITE, bold=True)
add_text(slide, 0.6, 2.5, 4.3, 0.35, "Bidirectional LSTM (hidden=128, 2-layer)",
         font_size=10, font_color=LIGHT_GRAY)

# 큰 F1
add_text(slide, 0.6, 3.1, 4.3, 0.4, "Weighted F1-Score",
         font_size=11, font_color=LIGHT_GRAY)
add_text(slide, 0.6, 3.45, 4.3, 1.4, "0.927",
         font_size=72, font_color=WHITE, bold=True)

# 작은 메트릭들
metrics_small = [("Acc", "0.928"), ("Prec", "0.928"), ("Rec", "0.928")]
for i, (k, v) in enumerate(metrics_small):
    x = 0.6 + i * 1.45
    add_text(slide, x, 5.05, 1.35, 0.3, k,
             font_size=10, font_color=LIGHT_GRAY)
    add_text(slide, x, 5.32, 1.35, 0.4, v,
             font_size=18, font_color=WHITE, bold=True)

add_text(slide, 0.6, 5.85, 4.3, 0.25,
         "Test set: 30,132 sequences (cycle-wise 20%)",
         font_size=8, font_color=LIGHT_GRAY)

# 우측: 3모델 최종 비교 표
add_text(slide, 5.1, 1.35, 7.8, 0.4, "3모델 최종 비교",
         font_size=14, font_color=NAVY, bold=True)

# 표 헤더
table_y2 = 1.85
add_rect(slide, 5.1, table_y2, 7.8, 0.4, fill=NAVY)
headers = ["Model", "Accuracy", "Precision", "F1"]
xs = [5.25, 7.85, 9.55, 11.25]
for h, x in zip(headers, xs):
    add_text(slide, x, table_y2 + 0.05, 1.8, 0.3, h,
             font_size=10, font_color=WHITE, bold=True)

# 표 row (최종 수치)
rows_data = [
    ("M1 · Sensor (BiLSTM)", "0.9276", "0.9276", "0.9267", BLUE, TEAL),
    ("M2 · Image (ResNet18)", "0.7263", "0.5276", "0.6112", TEAL, ORANGE),
    ("M3 · Fusion", "0.8941", "0.8891", "0.8824", CORAL, GREEN),
]
for i, (m, ac, pr, f1, mcolor, scolor) in enumerate(rows_data):
    y = table_y2 + 0.4 + i * 0.7
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, 5.1, y, 7.8, 0.7, fill=bg)
    add_circle(slide, 5.18, y + 0.25, 0.18, mcolor, "", font_size=8)
    add_text(slide, 5.4, y + 0.2, 2.35, 0.3, m,
             font_size=10, font_color=DARK_GRAY)
    add_text(slide, 7.75, y + 0.2, 1.7, 0.3, ac,
             font_size=11, font_color=NAVY, bold=True)
    add_text(slide, 9.45, y + 0.2, 1.7, 0.3, pr,
             font_size=11, font_color=NAVY, bold=True)
    add_text(slide, 11.15, y + 0.2, 1.6, 0.3, f1,
             font_size=11, font_color=scolor, bold=True)

# 우측 하단 박스: 핵심 관찰
add_rect(slide, 5.1, 4.7, 7.8, 1.45, fill=HIGHLIGHT)
add_text(slide, 5.25, 4.78, 7.5, 0.35, "핵심 관찰",
         font_size=12, font_color=NAVY, bold=True)
add_text(slide, 5.25, 5.15, 7.5, 0.95,
         "·  센서 단독(F1 0.927) > Fusion(F1 0.882) > 이미지 단독(F1 0.611)\n"
         "·  이미지 데이터 부족(9,542장, Cycle 4·9만)으로 Fusion에서 이미지 브랜치가 노이즈로 작용\n"
         "·  NSF-MAP P1 기본 Fusion 72%와 동일한 현상 — 데이터 균형이 Fusion 성능의 핵심",
         font_size=10, font_color=DARK_GRAY)


# ====================================================================
# Slide 10: 09 / 한계와 향후 개선 방향
# ====================================================================
print("[Slide 10] 09 한계와 개선")
slide = prs.slides[9]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "한계와 향후 개선 방향")

# 상단 요약: 현재 한계 진단
add_rect(slide, 0.42, 1.35, 12.5, 0.85, fill=SOFT_PINK)
add_text(slide, 0.6, 1.42, 12.2, 0.3, "진단: Fusion이 센서 단독보다 낮은 이유",
         font_size=12, font_color=CORAL, bold=True)
add_text(slide, 0.6, 1.72, 12.2, 0.45,
         "이미지 데이터 9,542장 (Cycle 4·9 한정, 전체의 ~18%) vs 센서 121,550 시퀀스 → "
         "Fusion에서 이미지 브랜치가 정보를 더하기보다 노이즈로 작용",
         font_size=10, font_color=DARK_GRAY)

# 좌측: 데이터 측면 개선
add_text(slide, 0.42, 2.45, 6.0, 0.4, "데이터 측면 개선",
         font_size=14, font_color=NAVY, bold=True)

data_improvements = [
    (BLUE, "1", "이미지 파트 확장",
     "Part 2~6 추가 사용 (최대 580GB)\n→ 9,542장 → 50K+ 장으로 ~5배 확장 기대"),
    (TEAL, "2", "Augmentation 고도화",
     "현재 Flip/Rotation/ColorJitter\n→ MixUp · CutMix · RandomErasing 추가"),
    (ORANGE, "3", "샘플링 균형 조정",
     "센서 시퀀스 다운샘플링 + 이미지 업샘플링\n→ Fusion 학습 시 두 모달리티 비율 맞춤"),
]
for i, (color, n, title, body) in enumerate(data_improvements):
    y = 2.95 + i * 1.05
    add_rect(slide, 0.42, y, 6.0, 0.92, fill=LIGHT_BG)
    add_circle(slide, 0.55, y + 0.22, 0.48, color, n, font_size=16)
    add_text(slide, 1.2, y + 0.08, 4.7, 0.35, title,
             font_size=12, font_color=NAVY, bold=True)
    add_text(slide, 1.2, y + 0.42, 4.7, 0.5, body,
             font_size=9, font_color=DARK_GRAY)

# 우측: 모델 측면 개선
add_text(slide, 6.92, 2.45, 6.0, 0.4, "모델·학습 측면 개선",
         font_size=14, font_color=NAVY, bold=True)

model_improvements = [
    (CORAL, "A", "Attention-based Fusion",
     "단순 Concat → Cross-Attention 도입\n→ 이미지 신뢰도가 낮을 때 센서에 가중치 부여"),
    (GREEN, "B", "Transfer Learning 강화",
     "ResNet18 backbone 동결 → 일부 unfreeze\n→ NSF-MAP P2 전략 (72%→88%) 차용"),
    (BLUE, "C", "Knowledge Infusion",
     "Process Ontology 기반 센서 범위 위반 패널티\n→ NSF-MAP P3 전략 (88%→93%) 차용"),
]
for i, (color, n, title, body) in enumerate(model_improvements):
    y = 2.95 + i * 1.05
    add_rect(slide, 6.92, y, 6.0, 0.92, fill=LIGHT_BG)
    add_circle(slide, 7.05, y + 0.22, 0.48, color, n, font_size=16)
    add_text(slide, 7.7, y + 0.08, 4.7, 0.35, title,
             font_size=12, font_color=NAVY, bold=True)
    add_text(slide, 7.7, y + 0.42, 4.7, 0.5, body,
             font_size=9, font_color=DARK_GRAY)

# 하단 박스
add_rect(slide, 0.42, 6.1, 12.5, 0.35, fill=NAVY)
add_text(slide, 0.6, 6.13, 12.2, 0.3,
         "→  데이터 균형 + Attention Fusion 적용 시 Fusion이 센서 단독을 상회할 것으로 기대 (NSF-MAP 사례 참조)",
         font_size=10, font_color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ====================================================================
# Slide 11: 10 / 스케줄 + 역할 분담
# ====================================================================
print("[Slide 11] 10 스케줄 + 역할")
slide = prs.slides[10]
remove_content_keep_header(slide, keep_indices={0, 1, 2})
set_header_title(slide, "스케줄 및 역할 분담")

# 좌측: 스케줄
add_text(slide, 0.42, 1.35, 6.0, 0.4, "주차별 진행 계획",
         font_size=14, font_color=NAVY, bold=True)

schedule = [
    ("1-2주차", "완료", "데이터셋 확보, EDA, 전처리 파이프라인", TEAL),
    ("3주차", "완료", "3개 모델 아키텍처 구현, Kaggle 환경 구축", TEAL),
    ("4주차", "완료", "센서·이미지·Fusion 3모델 학습 완료, 비교 분석", TEAL),
    ("5주차 (현재)", "진행중", "한계 분석 및 개선 방향 도출, 중간 발표 준비", ORANGE),
    ("6주차", "예정", "추가 개선 실험 (Attention Fusion 등), 최종 보고서", LIGHT_GRAY),
]
for i, (week, status, body, color) in enumerate(schedule):
    y = 1.85 + i * 0.78
    add_rect(slide, 0.42, y, 6.0, 0.68, fill=LIGHT_BG)
    # 좌측 컬러바
    add_rect(slide, 0.42, y, 0.15, 0.68, fill=color)
    add_text(slide, 0.7, y + 0.05, 1.6, 0.3, week,
             font_size=10, font_color=NAVY, bold=True)
    add_text(slide, 0.7, y + 0.34, 1.6, 0.3, status,
             font_size=9, font_color=color, bold=True)
    add_text(slide, 2.4, y + 0.15, 3.5, 0.5, body,
             font_size=10, font_color=DARK_GRAY,
             anchor=MSO_ANCHOR.MIDDLE)

# 우측: 역할 분담
add_text(slide, 6.92, 1.35, 6.0, 0.4, "역할 분담 (파트별)",
         font_size=14, font_color=NAVY, bold=True)

roles = [
    (BLUE, "하현성", [
        "멀티모달 융합 아키텍처 설계 및 코드 통합",
        "센서 브랜치 (BiLSTM) 구현 및 학습",
        "Kaggle 실행 환경 · 데이터 파이프라인",
    ]),
    (TEAL, "김주헌", [
        "이미지 브랜치 (ResNet18) 구현 및 학습",
        "데이터 EDA 및 라벨 정리 / 클래스 병합",
        "평가 지표 시각화 (Confusion Matrix 등)",
    ]),
    (CORAL, "공동 수행", [
        "Fusion 모델 통합 학습 및 디버깅",
        "3모델 성능 비교 분석",
        "발표 자료 및 기말 보고서 작성",
    ]),
]
y0 = 1.85
for i, (color, name, items) in enumerate(roles):
    y = y0 + i * 1.3
    add_rect(slide, 6.92, y, 6.0, 1.18, fill=LIGHT_BG)
    add_rect(slide, 6.92, y, 1.4, 1.18, fill=color)
    add_text(slide, 6.92, y + 0.35, 1.4, 0.5, name,
             font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    body_text = "\n".join(f"·  {it}" for it in items)
    add_text(slide, 8.45, y + 0.1, 4.4, 1.0, body_text,
             font_size=10, font_color=DARK_GRAY)


# ====================================================================
# Slide 12: Thank you (그대로)
# ====================================================================
print("[Slide 12] Thank you")
slide = prs.slides[11]
# GROUP 안의 "Thank you for your attention" 부분만 변경
group = list(slide.shapes)[0]
sub_shapes = [s for s in group.shapes if s.has_text_frame and s.text_frame.text.strip()]
new_thanks = ["Thank You", "Q & A"]
for i, sh in enumerate(sub_shapes):
    if i < len(new_thanks):
        tf = sh.text_frame
        p = tf.paragraphs[0]
        if p.runs:
            p.runs[0].text = new_thanks[i]
            for r in p.runs[1:]:
                r.text = ''
        else:
            p.text = new_thanks[i]
        for para in tf.paragraphs[1:]:
            para.text = ''


# ====================================================================
# 발표자 노트 초기화 (원본 논문 발표용 노트 제거)
# ====================================================================
print("\nClearing speaker notes...")
for slide in prs.slides:
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()


# ====================================================================
# 저장
# ====================================================================
print(f"\nSaving to: {OUT}")
prs.save(OUT)
print("Done.")
