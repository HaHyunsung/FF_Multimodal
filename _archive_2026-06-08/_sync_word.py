"""Word 문서를 PPT 노트 기준으로 동기화 (서식 보존)
- 슬라이드 2: 1 단락을 2 단락으로 분리
- 슬라이드 8 단락 3: 문구 교체
- 슬라이드 10 단락 1: 문구 교체
- 슬라이드 11 단락 1, 2: 문구 교체
- PPT 슬라이드 11의 "W평가" 오타도 함께 수정
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from docx import Document
from copy import deepcopy
from pptx import Presentation

DOCX = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표_스크립트.docx'
PPT  = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\3조_딥러닝 기반 제조 공정 이상 탐지_중간발표.pptx'

# ─── 1. PPT의 W평가 오타 먼저 수정 ────────────────────────
print("=== PPT 오타 수정 ===")
prs = Presentation(PPT)
slide11 = prs.slides[10]
ppt_changed = False
for shape in slide11.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if "W평가" in r.text:
                    print(f"  shape text: {r.text!r}")
                    r.text = r.text.replace("W평가", "평가")
                    print(f"  → {r.text!r}")
                    ppt_changed = True
# notes_slide도 확인
if slide11.has_notes_slide:
    notes_tf = slide11.notes_slide.notes_text_frame
    for p in notes_tf.paragraphs:
        for r in p.runs:
            if "W평가" in r.text:
                print(f"  notes: {r.text!r}")
                r.text = r.text.replace("W평가", "평가")
                ppt_changed = True
if ppt_changed:
    prs.save(PPT)
    print("PPT 저장됨.")
else:
    print("PPT에 W평가 미발견.")

# ─── 2. Word 문서 패치 ──────────────────────────────────────
print("\n=== Word 동기화 ===")
doc = Document(DOCX)

def replace_run_text(paragraph, new_text):
    """단락의 첫 run에 새 텍스트 넣고 나머지 run은 비움 (서식 유지)"""
    runs = paragraph.runs
    if not runs:
        paragraph.add_run(new_text)
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ''

def insert_paragraph_after(paragraph, text):
    """paragraph 바로 뒤에 동일한 서식의 새 단락 삽입"""
    new_p = deepcopy(paragraph._p)
    paragraph._p.addnext(new_p)
    # 새 단락의 모든 run 텍스트 초기화 + 첫 run에 text 넣기
    from docx.text.paragraph import Paragraph as DocParagraph
    wrapped = DocParagraph(new_p, paragraph._parent)
    runs = wrapped.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ''
    else:
        wrapped.add_run(text)
    return wrapped

# 슬라이드별 헤더 인덱스 매핑
slide_headers = {}
script_paras_by_slide = {}  # slide_num → [paragraph 객체 리스트]
current_slide = None
for idx, p in enumerate(doc.paragraphs):
    t = p.text.strip()
    if not t:
        continue
    if t.startswith("Slide ") and "—" in t:
        try:
            n = int(t.split("Slide ")[1].split(" —")[0])
            current_slide = n
            slide_headers[n] = idx
            script_paras_by_slide[n] = []
        except:
            current_slide = None
    elif t.startswith("PART 2") or "Q&A 대응" in t:
        current_slide = None
    elif current_slide is not None:
        script_paras_by_slide[current_slide].append(p)

# ── Slide 2: 1단락을 2단락으로 분리 ──
print("\n[Slide 2] 1단락 → 2단락 분리")
slide2_paras = script_paras_by_slide[2]
if len(slide2_paras) == 2:
    # 단락 1을 PPT 단락 1로 교체
    p1_new = "프로젝트의 출발점은 저희 팀의 현장 경험입니다. 팀원 전원이 자동화 장비 업계에서 근무하면서 제조 공정에서 딥러닝이 어떻게 활용되는지 직접 접하고 있습니다."
    p2_insert = "현재 산업 현장의 불량 판정은 대부분 비전 카메라 단독 방식입니다. 그러나 실제로는 조명, 각도, 가려짐 때문에 비전만으로 판단이 어려운 경우가 자주 발생합니다. 반대로 센서 데이터도 노이즈나 고장 때문에 단독으로는 신뢰하기 어려운 경우가 존재합니다."
    # 기존 단락 1은 PPT 단락 1로
    replace_run_text(slide2_paras[0], p1_new)
    # 단락 1 뒤에 PPT 단락 2 삽입
    insert_paragraph_after(slide2_paras[0], p2_insert)
    # 기존 단락 2는 그대로 (PPT 단락 3과 동일)
    print(f"  완료")

# ── Slide 8 단락 3 교체 ──
print("\n[Slide 8] 단락 3 교체")
slide8_paras = script_paras_by_slide[8]
if len(slide8_paras) >= 3:
    new_text = "이 문제를 Kaggle Notebook으로 해결했습니다. Kaggle 서버에 마운트된 원본 데이터를 다운로드 없이 직접 읽고, 무료로 제공되는 Tesla T4 GPU로 실험 사이클을 단축했습니다. 검증된 모델은 추후 로컬에서도 재현도 고려 중입니다."
    replace_run_text(slide8_paras[2], new_text)
    print(f"  완료")

# ── Slide 10 단락 1 교체 ──
print("\n[Slide 10] 단락 1 교체")
slide10_paras = script_paras_by_slide[10]
if len(slide10_paras) >= 1:
    new_text = "한계의 핵심은 데이터 불균형입니다. 이미지 9,500장 대비 센서 12만 시퀀스로 Fusion에서 이미지 브랜치가 노이즈로 작용했을 것으로 추측됩니다."
    replace_run_text(slide10_paras[0], new_text)
    print(f"  완료")

# ── Slide 11 단락 1, 2 교체 ──
print("\n[Slide 11] 단락 1, 2 교체")
slide11_paras = script_paras_by_slide[11]
if len(slide11_paras) >= 2:
    p1_new = "스케줄은 1-2주차 데이터 확보와 모델 구현, 3-4주차에 3개 모델 학습과 비교 분석까지 진행했습니다. 남은 5-6주차는 추가 개선 실험과 최종 보고서 작성입니다."
    # "W평가"는 사용자 결정에 따라 "평가"로 수정
    p2_new = "역할은 하현성이 멀티모달 융합 아키텍처와 센서 브랜치, Kaggle 환경 구축을, 김주헌이 이미지 브랜치와 데이터 탐색 및 라벨 정리, 평가 시각화를 담당했습니다. Fusion 통합과 비교 분석, 보고서 작성은 공동 수행입니다."
    replace_run_text(slide11_paras[0], p1_new)
    replace_run_text(slide11_paras[1], p2_new)
    print(f"  완료")

doc.save(DOCX)
print("\nWord 저장 완료.")
