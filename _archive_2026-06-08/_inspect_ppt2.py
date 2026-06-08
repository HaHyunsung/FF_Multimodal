"""GROUP 포함 전체 텍스트 인벤토리 추출 - 슬라이드 수정용 매핑 작성"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_text_frames(shape, path=""):
    """모든 텍스트 프레임을 (path, shape) 형태로 yield"""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for i, sub in enumerate(shape.shapes):
            yield from iter_text_frames(sub, f"{path}.{i}")
    elif shape.has_text_frame and shape.text_frame.text.strip():
        yield path, shape


p = Presentation(r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx')

for i, slide in enumerate(p.slides):
    print(f"\n=== Slide {i+1} ===")
    for j, shape in enumerate(slide.shapes):
        for path, s in iter_text_frames(shape, str(j)):
            txt = s.text_frame.text.replace('\n', ' | ')
            print(f"  [{path}] \"{txt}\"")
