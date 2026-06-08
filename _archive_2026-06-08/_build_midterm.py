"""
중간발표 PPT 생성:
- 기존 딥러닝응용_3조_중간발표.pptx 템플릿을 그대로 유지
- 각 슬라이드 텍스트 박스의 내용만 교체
- 폰트/색상/레이아웃은 run 단위 텍스트 교체로 자동 유지
"""
import sys, io, shutil, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'
# 원본 백업하고 새로 작성
OUT = SRC


# ──────────────────────────────────────────────
# 슬라이드별 텍스트 교체 매핑
# 키: shape path (점 표기)
# 값: 새 텍스트 (여러 줄은 \n)
# ──────────────────────────────────────────────

SLIDE_CONTENT = {
    # === Slide 1: 표지 ===
    1: {
        "1.0": "멀티모달 딥러닝 기반\n제조 공정 이상 탐지\n(이미지 + 시계열 센서)",
        "1.1": "Future Factories Dataset 활용\n이미지 단독 · 센서 단독 · 멀티모달 융합 모델 비교",
        "1.2": "아주대학교 딥러닝응용 기말 프로젝트",
        "1.3": "[3조] 김주헌, 하현성",
    },

    # === Slide 2: 서론 (문제 제기 → 프로젝트 동기) ===
    2: {
        "2": "서론: 프로젝트 동기",
        "5": "비전 단독 판정의 한계",
        "6": "현업 불량 판정 대부분이 비전 카메라 단독 방식\n카메라 각도·조명·가려짐으로 인한 오판정 빈번",
        "9": "센서 단독의 신뢰성 문제",
        "10": "센서 노이즈·고장으로 단독 신뢰가 어려운 상황 존재\n두 모달리티가 서로 보완할 가능성",
        "13": "팀의 현장 경험",
        "14": "팀원 전원 자동화 장비 업계 종사\n제조 공정·센서 데이터에 대한 실무 이해 보유",
        "16": "→  이미지와 센서를 함께 활용하는 멀티모달 접근을 직접 구현하여, 단일 모달리티 대비 실제 성능 차이를 정량적으로 확인",
    },

    # === Slide 3: 프로젝트 목표 ===
    3: {
        "2": "프로젝트 목표와 접근",
        "3.0": "3-Model Comparison Framework",
        "3.3": "센서 단독 모델",
        "3.4": "BiLSTM 기반\n시계열 센서 22채널\n이상 유형 분류",
        "3.7": "이미지 단독 모델",
        "3.8": "ResNet18 (Transfer Learning)\n카메라 이미지 224×224\n이상 유형 분류",
        "3.11": "멀티모달 융합 모델",
        "3.12": "Decision-Level Fusion\n이미지 + 센서 특징 결합\n상호 보완 효과 검증",
        "3.14": "세 모델을 동일한 데이터·평가 지표로 비교하여, 멀티모달 융합이 단일 모달리티 대비 실제로 어떤 성능 차이를 만드는지 정량 분석",
    },

    # === Slide 4: 관련 연구 (NSF-MAP) ===
    4: {
        "2": "관련 연구: NSF-MAP (IJCAI 2025)",
        "3": "Decision-level Fusion",
        "5": "• EfficientNet-B0 + Autoencoder\n• 특징 벡터 Concatenate\n• 시계열·이미지 동시 활용",
        "7": "기본 Fusion 구조에서\nAccuracy 72% 달성",
        "8": "Transfer Learning",
        "10": "• Encoder 사전학습 후 동결\n• Decoder만 Fine-tuning\n• 과적합 방지, 안정적 특징 추출",
        "12": "Fusion 효과 발휘\nAccuracy 88%로 향상",
        "13": "Knowledge Infusion",
        "15": "• Process Ontology 활용\n• 센서 정상 범위 위반 시 패널티\n• 도메인 지식이 학습 가이드",
        "17": "Best 성능\nAccuracy 93%, F1 93%",
        "19": "→  본 프로젝트는 NSF-MAP의 단계별 접근을 참고하되, 강의에서 학습한 LSTM·ResNet으로 재구성하여 구현",
    },

    # === Slide 5: 데이터셋 ===
    5: {
        "2": "데이터셋: Future Factories",
        "3": "Future Factories Lab Dataset (USC)",
        "5": "조립 라인:  산업용 로봇 4대 + 컨베이어",
        "7": "조립 대상:  4개 부품 로켓 모형 (의도적 결함 주입)",
        "9": "수집 규모:  30시간 연속, 166,001 레코드, 285 사이클",
        "11": "센서 데이터:  그리퍼 로드셀·포텐셔미터, 로봇 관절각도, 컨베이어 온도 등 40+ 채널",
        "13": "이미지 데이터:  카메라 2대 동기화 촬영 (약 166K 장)",
        "14": "라벨 분포 (소수 클래스 병합 후)",
        "17": "데이터 분할 전략",
        "18": "Cycle-wise split (Train 80% / Test 20%) — 같은 사이클 데이터가 train/test에 섞이지 않도록 처리 (Data leakage 방지)",
    },

    # === Slide 6: 우리 아키텍처 ===
    6: {
        "2": "전체 아키텍처",
        "3.0": "Sensor Data\n(22 channels)",
        "3.2": "BiLSTM\n(hidden 128, 2-layer)",
        "3.4": "Sensor\nFeatures (256)",
        "3.5": "Image Data\n(224×224, Cycle 4·9)",
        "3.7": "ResNet18",
        "3.9": "(ImageNet\npretrained)",
        "3.11": "Image\nFeatures (512)",
        "3.12": "Concat\n[fS ; fI]\n(768)",
        "3.14": "FC + Dropout\n+ Classifier",
        "3.15": "Anomaly\nType (5-class)",
        "3.19": "센서 브랜치 (전체 시점)",
        "3.20": "22채널 센서 시퀀스를 BiLSTM에 입력\n마지막 hidden state로 256차원 특징 벡터 추출",
        "3.22": "이미지 브랜치 (Cycle 4·9 시점)",
        "3.23": "ResNet18 사전학습 가중치로 512차원 특징 추출\n분류 레이어 제거 후 backbone 동결",
    },

    # === Slide 7: 모델별 학습 전략 ===
    7: {
        "2": "모델별 학습 전략",
        "3.1": "M1",
        "3.2": "Sensor-Only (BiLSTM)",
        "3.3": "• 시계열 시퀀스 (T=50, 22ch)\n• Sliding Window\n• Cycle-wise 분할\n• 5-class 분류",
        "3.4": "전체 사이클 활용",
        "3.7": "M2",
        "3.8": "Image-Only (ResNet18)",
        "3.9": "• ImageNet 사전학습 활용\n• Data Augmentation\n• Fine-tuning\n• Cycle 4·9 이미지만 사용",
        "3.10": "Transfer Learning",
        "3.13": "M3",
        "3.14": "Multimodal Fusion",
        "3.15": "• Decision-Level Fusion\n• 이미지·센서 특징 Concat\n• FC Layer로 분류\n• 두 모달리티 상호 보완",
        "3.16": "Sensor + Image",
        "3.17": "공통 전처리·평가 지표 사용 → 세 모델의 효과를 동일 조건에서 비교",
    },

    # === Slide 8: 실행 환경 (Kaggle) ===
    8: {
        "2": "실행 환경 및 데이터 제약",
        "3.0": "이미지 데이터 용량 문제와 해결",
        "3.2": "FF Dataset 원본 이미지 용량",
        "3.4": "Kaggle 분할 다운로드 (6 파트)",
        "3.6": "Kaggle Notebook 활용",
        "3.8": "무료 GPU (Tesla T4) 제공",
        "3.10": "GitHub로 코드 관리 및 협업",
        "3.11": "프로젝트 인프라 구성",
        "3.13": "1",
        "3.14": "데이터 접근 전략",
        "3.15": "이미지 원본 580GB → 로컬 다운로드 불가\nKaggle 서버 상의 데이터셋을 Notebook에 마운트하여 직접 사용\n→ 데이터 다운로드 없이 학습 가능",
        "3.17": "2",
        "3.18": "이미지 샘플링",
        "3.19": "6 파트 중 1 파트만 사용\n전체 53,881장 중 9,542장 매칭 (약 18%)\n→ 학술적으로 충분한 샘플 (NSF-MAP은 약 15K장 사용)",
        "3.21": "3",
        "3.22": "재현 가능성 확보",
        "3.23": "GitHub Repo: github.com/HaHyunsung/FF_Multimodal\n노트북·전처리 코드·모델 정의 모두 공개\n팀원 간 코드 동기화 및 버전 관리",
        "3.25": "GPU 환경",
        "3.26": "Tesla T4 × 2 (Kaggle 무료 할당, 주 30시간)\n→ ResNet18 학습 시간 단축, 멀티모달 실험 가능",
    },

    # === Slide 9: 중간 결과 ===
    9: {
        "2": "중간 결과: 센서 단독 모델 (완료)",
        "4.0": "* Model 2 (이미지), Model 3 (Fusion)은 학습 진행 중 — 최종 발표에서 보완",
        "4.1": "Test 세트 평가 (Cycle-wise 20% split)",
        "4.3": "주요 관찰 사항",
        "4.4": "BiLSTM이 시계열 패턴만으로도 0.93 수준의 안정적 분류 성능을 보여, 센서 데이터의 정보량이 풍부함을 확인",
    },

    # === Slide 10: 진행 현황 + 향후 계획 ===
    10: {
        "2": "진행 현황과 향후 계획",
        "3.0": "Model 1 (Sensor)",
        "3.3": "완료",
        "3.4": "Model 2 (Image)",
        "3.7": "학습 중",
        "3.8": "Model 3 (Fusion)",
        "3.11": "대기",
        "3.12": "추가 실험",
        "3.15": "계획",
        "3.16": "F1 0.927",
        "3.17": "BATCH 매칭",
        "3.18": "구조 완성",
        "3.20": "남은 작업 (5~6주차)",
        "3.21": "Model 2·3 학습 완료 후 동일 지표로 비교\n→ Confusion Matrix, 학습 곡선, 클래스별 F1 분석",
        "3.23": "추가 분석 방향",
        "3.24": "(1) Fusion 전략 비교: Concat vs Attention\n(2) Data Augmentation 효과 검증\n(3) 클래스 불균형 처리 비교 (Class Weight vs Focal Loss)",
    },

    # === Slide 11: 스케줄 + 역할 분담 ===
    11: {
        "2": "스케줄과 역할 분담",
        "3.0": "주차별 진행 계획",
        "3.2": "1-3주차 ✓ : 데이터셋 확보, 전처리 파이프라인, 모델 아키텍처 구현 완료",
        "3.4": "4주차 (현재) ✓ : 센서 모델 학습 완료(F1 0.927), 이미지·Fusion 학습 진행 중",
        "3.6": "5주차 : 3모델 비교 실험 완료, 추가 Fusion 전략 실험",
        "3.8": "6주차 : 최종 발표 자료 및 보고서 작성, 결과 정리",
        "3.9": "역할 분담 (파트별)",
        "3.11": "[하현성] 멀티모달 융합 아키텍처 설계 및 코드 통합 — Decision-Level Fusion 구조, Kaggle 실행 환경 구축, GitHub 코드 관리",
        "3.13": "[하현성] 센서 브랜치 구현 — BiLSTM 모델 정의, 시계열 시퀀스 전처리, 클래스 불균형 처리",
        "3.15": "[김주헌] 이미지 브랜치 구현 — ResNet18 Transfer Learning, 이미지 전처리 및 Augmentation 파이프라인",
        "3.17": "[김주헌] 데이터 분석 및 평가 — EDA, 라벨 정리, 평가 지표 설계, 시각화 (Confusion Matrix, 학습 곡선)",
        "3.19": "공동 수행",
        "3.20": "Fusion 모델 통합 학습 및 디버깅  ·  3모델 성능 비교 분석  ·  발표 자료 및 기말 보고서 작성",
    },

    # === Slide 12: Thank you ===
    12: {
        "0.0": "Thank you",
        "0.1": "Q & A",
    },
}


# ──────────────────────────────────────────────
# 텍스트 교체 헬퍼 (run 단위로 첫 번째 run에 모든 텍스트 몰아넣고
#  나머지 run은 비워서 폰트 스타일 유지)
# ──────────────────────────────────────────────
def set_text_preserve_style(text_frame, new_text):
    """텍스트 프레임의 텍스트를 교체하되 첫 paragraph의 첫 run의 스타일을 유지한다."""
    lines = new_text.split('\n')

    # 첫 paragraph 확인
    if not text_frame.paragraphs:
        return

    first_para = text_frame.paragraphs[0]

    # 첫 paragraph의 첫 run 보존 (스타일 복사용)
    if not first_para.runs:
        # run이 없으면 paragraph 텍스트로만 설정
        first_para.text = lines[0]
    else:
        first_run = first_para.runs[0]
        first_run.text = lines[0]
        # 첫 run을 제외한 나머지 run 텍스트는 비우기
        for run in first_para.runs[1:]:
            run.text = ''

    # 나머지 라인을 추가 paragraph로
    # 기존 paragraph가 더 있으면 활용, 부족하면 새로 생성
    existing_paras = list(text_frame.paragraphs)

    for i, line in enumerate(lines[1:], start=1):
        if i < len(existing_paras):
            para = existing_paras[i]
            if para.runs:
                para.runs[0].text = line
                for run in para.runs[1:]:
                    run.text = ''
            else:
                para.text = line
        else:
            # 새 paragraph (이전 paragraph 스타일 따라감)
            new_p = copy.deepcopy(first_para._p)
            # text 비우기
            from pptx.oxml.ns import qn
            for r in new_p.findall(qn('a:r')):
                t = r.find(qn('a:t'))
                if t is not None:
                    t.text = ''
            # 첫 run만 남기고 나머지 제거
            runs = new_p.findall(qn('a:r'))
            for r in runs[1:]:
                new_p.remove(r)
            # 첫 run 텍스트 설정
            r = new_p.find(qn('a:r'))
            if r is not None:
                t = r.find(qn('a:t'))
                if t is not None:
                    t.text = line
            text_frame._txBody.append(new_p)

    # 사용하지 않는 추가 paragraph 비우기
    if len(lines) < len(existing_paras):
        for j in range(len(lines), len(existing_paras)):
            existing_paras[j].text = ''


def find_shape_by_path(shapes, path):
    """점 표기 경로로 shape를 찾는다. 예: '3.0.1'"""
    parts = path.split('.')
    current = shapes[int(parts[0])]
    for p in parts[1:]:
        current = current.shapes[int(p)]
    return current


# ──────────────────────────────────────────────
# 실행
# ──────────────────────────────────────────────
print(f"Loading: {SRC}")
prs = Presentation(SRC)

modified_count = 0
for slide_idx, content_map in SLIDE_CONTENT.items():
    slide = prs.slides[slide_idx - 1]
    print(f"\nSlide {slide_idx}: {len(content_map)} text replacements")
    for path, new_text in content_map.items():
        try:
            shape = find_shape_by_path(slide.shapes, path)
            if shape.has_text_frame:
                set_text_preserve_style(shape.text_frame, new_text)
                modified_count += 1
            else:
                print(f"  [WARN] {path}: shape has no text frame")
        except Exception as e:
            print(f"  [ERROR] {path}: {e}")

print(f"\nTotal replacements: {modified_count}")
print(f"Saving to: {OUT}")
prs.save(OUT)
print("Done.")
