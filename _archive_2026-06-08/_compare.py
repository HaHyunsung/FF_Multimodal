"""PPT 노트와 Word 스크립트 비교"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from docx import Document
from pptx import Presentation

PPT = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\3조_딥러닝 기반 제조 공정 이상 탐지_중간발표.pptx'
DOCX = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표_스크립트.docx'

# PPT 노트
prs = Presentation(PPT)
ppt_notes = {}
for i, slide in enumerate(prs.slides, 1):
    if slide.has_notes_slide:
        t = slide.notes_slide.notes_text_frame.text.strip()
        if t:
            ppt_notes[i] = t

# Word 본문 추출 - 슬라이드별로 그룹화
doc = Document(DOCX)
print("=== Word 문서 단락 구조 ===")
current_slide = None
word_scripts = {}
for idx, p in enumerate(doc.paragraphs):
    t = p.text.strip()
    if not t:
        continue
    # "Slide N —" 헤더 감지
    if t.startswith("Slide ") and "—" in t:
        try:
            n = int(t.split("Slide ")[1].split(" —")[0])
            current_slide = n
            word_scripts[n] = []
            print(f"\n[Slide {n}] {t}")
        except:
            pass
    elif t.startswith("PART 2") or t.startswith("PART 1"):
        if "PART 2" in t:
            current_slide = None
        print(f"\n>>> {t}")
    elif current_slide is not None:
        word_scripts[current_slide].append(t)
        print(f"  · {t[:80]}")

print("\n\n=== PPT 노트 vs Word 비교 ===")
for slide_num in sorted(ppt_notes.keys()):
    ppt_text = ppt_notes[slide_num]
    word_paras = word_scripts.get(slide_num, [])
    word_combined = "\n".join(word_paras)

    # PPT 노트를 \n\n 단위로 분리 후 비교
    ppt_paras = [p.strip() for p in ppt_text.split('\n') if p.strip()]

    if len(ppt_paras) != len(word_paras):
        print(f"\n[Slide {slide_num}] 단락 수 다름: PPT {len(ppt_paras)} vs Word {len(word_paras)}")

    # 각 단락 비교
    for i, ppt_p in enumerate(ppt_paras):
        word_p = word_paras[i] if i < len(word_paras) else "(없음)"
        if ppt_p != word_p:
            print(f"\n[Slide {slide_num}] 단락 {i+1} 차이:")
            print(f"  PPT:  {ppt_p[:150]}")
            print(f"  Word: {word_p[:150]}")
