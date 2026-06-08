"""Slide 11 김주헌 역할 텍스트의 'EDA' 표현 교체"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'

prs = Presentation(SRC)
slide = prs.slides[10]  # Slide 11

OLD_TEXT = "데이터 EDA 및 라벨 정리 / 클래스 병합"
NEW_TEXT = "데이터 탐색 및 라벨 정리 / 클래스 병합"

found = False
for shape in slide.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if "EDA" in r.text:
                    print(f"Before: {r.text!r}")
                    r.text = r.text.replace("데이터 EDA", "데이터 탐색")
                    r.text = r.text.replace(", EDA,", ", 데이터 탐색,")
                    print(f"After:  {r.text!r}")
                    found = True

if not found:
    print("EDA 텍스트를 직접 못 찾음. 단락 전체 텍스트 검색...")
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "EDA" in t:
                print(f"Found in shape: {t[:100]}")

prs.save(SRC)
print("저장 완료.")
