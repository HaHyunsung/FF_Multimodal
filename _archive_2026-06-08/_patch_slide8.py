"""Slide 8의 '로컬 GPU 부재' 항목만 교체"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'

# 새 내용
NEW_TITLE = "빠른 프로토타이핑 필요"
NEW_BODY = ("3개 모델 × 학습 반복 시 로컬 자원 장시간 점유\n"
            "→ Kaggle T4로 실험 사이클 단축, 추후 로컬 재현 가능")

prs = Presentation(SRC)
slide = prs.slides[7]  # Slide 8 (0-index)

# "로컬 GPU 부재" 텍스트가 있는 도형 찾기 → 그 다음 본문 도형도 찾기
title_shape = None
body_shape = None
for shape in slide.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == "로컬 GPU 부재":
            title_shape = shape
        elif "팀원 PC에 학습용 GPU" in t:
            body_shape = shape

if title_shape and body_shape:
    # title 교체 (run 단위로 스타일 유지)
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = NEW_TITLE
        for r in p.runs[1:]:
            r.text = ''
    else:
        p.text = NEW_TITLE
    for para in tf.paragraphs[1:]:
        para.text = ''

    # body 교체
    tf = body_shape.text_frame
    lines = NEW_BODY.split('\n')
    # 첫 paragraph
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r.text = ''
    else:
        p0.text = lines[0]
    # 두 번째 paragraph
    existing = list(tf.paragraphs)
    if len(existing) >= 2:
        p1 = existing[1]
        if p1.runs:
            p1.runs[0].text = lines[1] if len(lines) > 1 else ''
            for r in p1.runs[1:]:
                r.text = ''
        else:
            p1.text = lines[1] if len(lines) > 1 else ''
        # 나머지 비우기
        for para in existing[2:]:
            para.text = ''
    else:
        # 두 번째 paragraph 없으면 추가
        if len(lines) > 1:
            np = tf.add_paragraph()
            np.text = lines[1]

    print(f"교체 완료: '{NEW_TITLE}'")
else:
    print(f"도형 못 찾음. title={title_shape}, body={body_shape}")

prs.save(SRC)
print("저장 완료.")
