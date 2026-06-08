import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
p = Presentation(r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx')
print(f'Slides: {len(p.slides)}')
for i, slide in enumerate(p.slides):
    print(f'--- Slide {i+1} (layout: {slide.slide_layout.name}) ---')
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            txt = shape.text_frame.text[:100].replace('\n', ' | ')
            print(f'  shape[{j}] {shape.shape_type} \"{txt}\"')
        else:
            print(f'  shape[{j}] {shape.shape_type}')
