"""PPT의 'W평가' 오타 수정 (run 분리 케이스 대응)"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

PPT = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\3조_딥러닝 기반 제조 공정 이상 탐지_중간발표.pptx'

prs = Presentation(PPT)
slide11 = prs.slides[10]

def fix_paragraph(p, label=""):
    """단락 전체 텍스트를 합쳐서 W평가 찾고, 발견 시 첫 run에 전체 텍스트 재배치"""
    full = "".join(r.text for r in p.runs)
    if "W평가" in full:
        new_full = full.replace("W평가", "평가")
        print(f"  [{label}] FOUND: {full[:100]!r}")
        # 첫 run에 전체 텍스트 넣고 나머지 비움
        if p.runs:
            p.runs[0].text = new_full
            for r in p.runs[1:]:
                r.text = ''
        print(f"  → fixed")
        return True
    return False

print("=== 슬라이드 11 shape 검사 ===")
for si, shape in enumerate(slide11.shapes):
    if shape.has_text_frame:
        for pi, p in enumerate(shape.text_frame.paragraphs):
            fix_paragraph(p, f"shape{si} para{pi}")

print("\n=== 슬라이드 11 notes 검사 ===")
if slide11.has_notes_slide:
    notes_tf = slide11.notes_slide.notes_text_frame
    for pi, p in enumerate(notes_tf.paragraphs):
        fix_paragraph(p, f"notes para{pi}")

prs.save(PPT)
print("\nPPT 저장 완료.")
