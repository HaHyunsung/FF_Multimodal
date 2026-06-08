"""Slide 8의 발표자 노트 교체"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\Hyunsung Ha\Desktop\아주대학교\4학년\딥러닝응용\기말 멀티모달 프로젝트\딥러닝응용_3조_중간발표.pptx'

NEW_NOTE = (
    "실행 환경에서 가장 큰 제약은 데이터 용량이었습니다. 이미지 원본이 6 파트 580GB로 로컬 PC 다운로드가 사실상 불가능했습니다.\n\n"
    "또한 3개 모델을 반복적으로 학습·실험해야 하는 상황에서 로컬 자원을 장시간 점유하는 부담이 컸기 때문에, 빠른 프로토타이핑이 필요했습니다.\n\n"
    "이 문제를 Kaggle Notebook으로 해결했습니다. Kaggle 서버에 마운트된 원본 데이터를 다운로드 없이 직접 읽고, 무료로 제공되는 Tesla T4 GPU로 실험 사이클을 단축했습니다. 검증된 모델은 추후 로컬에서도 재현 가능합니다.\n\n"
    "다만 이미지는 6 파트 중 1 파트만 사용해서 약 9,500장이 매칭됐는데, 이는 NSF-MAP에서 사용한 약 1.5만 장과 비교해도 학술적으로 유의미한 규모입니다.\n\n"
    "추가로 Kaggle 세션이 idle 시 자동 종료되는 점을 고려해서 Random Seed를 42로 고정하고, 학습된 가중치와 결과를 Kaggle Output에 별도 저장해 재현성을 확보했습니다."
)

prs = Presentation(SRC)
slide = prs.slides[7]  # Slide 8

notes_tf = slide.notes_slide.notes_text_frame
notes_tf.clear()

paragraphs = NEW_NOTE.split('\n\n')
for j, para_text in enumerate(paragraphs):
    if j == 0:
        p = notes_tf.paragraphs[0]
    else:
        p = notes_tf.add_paragraph()
    p.text = para_text

prs.save(SRC)
print(f"Slide 8 노트 교체 완료 ({len(NEW_NOTE)} chars)")
